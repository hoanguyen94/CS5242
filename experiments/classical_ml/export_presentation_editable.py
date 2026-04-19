"""
Export presentation as a native, editable PPTX using python-pptx.
All text is editable; images are embedded at full resolution.
"""

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DIR = Path(__file__).parent
IMG_DIR = DIR  # PNGs live alongside the script
OUT = DIR / f"presentation_{datetime.now().strftime('%y%m%d%H%M%S')}.pptx"

# ── Colours ──────────────────────────────────────────────────────
BLUE       = RGBColor(0x15, 0x65, 0xC0)
MID_BLUE   = RGBColor(0x19, 0x76, 0xD2)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
RED        = RGBColor(0xC6, 0x28, 0x28)

# ── Slide dimensions (16:9) ──────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use the blank layout for all slides
BLANK_LAYOUT = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────
def _add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _set_run(run, text, size=18, bold=False, italic=False, color=DARK_GREY, font_name="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name


def _add_para(tf, text, size=18, bold=False, italic=False, color=DARK_GREY,
              alignment=PP_ALIGN.LEFT, space_after=Pt(6), level=0, bullet=False):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = space_after
    p.level = level
    _set_run(p.runs[0] if p.runs else p.add_run(), text, size, bold, italic, color)
    if not bullet:
        p._pPr.attrib.pop('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone', None)
        # Remove bullet
        from pptx.oxml.ns import qn
        buNone = p._pPr.makeelement(qn('a:buNone'), {})
        p._pPr.append(buNone)
    return p


def _first_para(tf, text, size=18, bold=False, italic=False, color=DARK_GREY, alignment=PP_ALIGN.LEFT):
    """Set text on the first (existing) paragraph of a textframe."""
    p = tf.paragraphs[0]
    p.alignment = alignment
    _set_run(p.runs[0] if p.runs else p.add_run(), text, size, bold, italic, color=color)
    return p


def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    # Title
    txBox = _add_textbox(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    _first_para(tf, title, size=40, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    # Subtitle
    if subtitle:
        _add_para(tf, subtitle, size=24, color=MID_BLUE, alignment=PP_ALIGN.CENTER, space_after=Pt(4))
    return slide


def add_content_slide(title, bullets=None, notes=""):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    # Title bar
    txBox = _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = txBox.text_frame; tf.word_wrap = True
    _first_para(tf, title, size=32, bold=True, color=BLUE)
    # Body
    if bullets:
        body = _add_textbox(slide, Inches(0.7), Inches(1.3), Inches(11.9), Inches(5.5))
        bf = body.text_frame; bf.word_wrap = True
        for i, b in enumerate(bullets):
            if i == 0:
                _first_para(bf, b, size=18, color=DARK_GREY)
            else:
                _add_para(bf, b, size=18)
    return slide


def add_image_slide(title, img_filename, caption=""):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    # Title
    txBox = _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = txBox.text_frame; tf.word_wrap = True
    _first_para(tf, title, size=32, bold=True, color=BLUE)
    # Image – centred
    img_path = IMG_DIR / img_filename
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(1.5), Inches(1.3), Inches(10.3), Inches(5.3))
    else:
        b = _add_textbox(slide, Inches(2), Inches(3), Inches(9), Inches(1))
        _first_para(b.text_frame, f"[Image not found: {img_filename}]", size=20, color=RED,
                    alignment=PP_ALIGN.CENTER)
    # Caption
    if caption:
        c = _add_textbox(slide, Inches(0.7), Inches(6.7), Inches(11.9), Inches(0.6))
        cf = c.text_frame; cf.word_wrap = True
        _first_para(cf, caption, size=14, italic=True, color=MID_BLUE, alignment=PP_ALIGN.CENTER)
    return slide


def add_table_slide(title, headers, rows, col_widths=None, intro=""):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    # Title
    txBox = _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = txBox.text_frame; tf.word_wrap = True
    _first_para(tf, title, size=32, bold=True, color=BLUE)

    top = Inches(1.3)
    if intro:
        ib = _add_textbox(slide, Inches(0.7), Inches(1.2), Inches(11.9), Inches(0.6))
        ibf = ib.text_frame; ibf.word_wrap = True
        _first_para(ibf, intro, size=16, color=DARK_GREY)
        top = Inches(1.9)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_w = Inches(11.5) if col_widths is None else sum(col_widths)
    tbl_h = Inches(0.4 * n_rows)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.9), top, tbl_w, tbl_h)
    table = shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = MID_BLUE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GREY
                p.alignment = PP_ALIGN.CENTER

    return slide


def add_two_col_slide(title, left_title, left_bullets, right_title, right_bullets, footer=""):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    # Title
    txBox = _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = txBox.text_frame; tf.word_wrap = True
    _first_para(tf, title, size=32, bold=True, color=BLUE)

    for col_idx, (ct, cb) in enumerate([(left_title, left_bullets), (right_title, right_bullets)]):
        left = Inches(0.7) if col_idx == 0 else Inches(6.8)
        # Column title
        ct_box = _add_textbox(slide, left, Inches(1.3), Inches(5.8), Inches(0.5))
        _first_para(ct_box.text_frame, ct, size=22, bold=True, color=MID_BLUE)
        # Column bullets
        cb_box = _add_textbox(slide, left, Inches(1.9), Inches(5.8), Inches(4.5))
        cbf = cb_box.text_frame; cbf.word_wrap = True
        for i, b in enumerate(cb):
            if i == 0:
                _first_para(cbf, f"• {b}", size=16, color=DARK_GREY)
            else:
                _add_para(cbf, f"• {b}", size=16)

    if footer:
        fb = _add_textbox(slide, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.7))
        fbf = fb.text_frame; fbf.word_wrap = True
        _first_para(fbf, footer, size=14, italic=True, color=MID_BLUE)

    return slide


def add_quote_slide(title, quotes):
    """Add a slide with a title and blockquote-style text."""
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    txBox = _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = txBox.text_frame; tf.word_wrap = True
    _first_para(tf, title, size=32, bold=True, color=BLUE)

    body = _add_textbox(slide, Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.0))
    bf = body.text_frame; bf.word_wrap = True
    for i, q in enumerate(quotes):
        if i == 0:
            _first_para(bf, q, size=18, color=DARK_GREY)
        else:
            _add_para(bf, q, size=18)
    return slide


# ═════════════════════════════════════════════════════════════════
# SLIDES
# ═════════════════════════════════════════════════════════════════

# 1 ── Title
add_title_slide(
    "Traditional ML on Mini-ImageNet",
    "Pretrained Backbones as Frozen Feature Extractors — CS5242 Project, Section 2",
)

# 2 ── Agenda
add_content_slide("Agenda", [
    "1. Problem & Motivation",
    "2. Approach: Frozen Backbone + Classical Classifier",
    "3. Experimental Design",
    "4. Results: Backbone & Classifier Comparison (32×32)",
    "5. Results: Resolution Impact (224×224 vs 32×32)",
    "6. t-SNE Feature Visualisation",
    "7. Key Findings & Backbone Downselection",
    "8. Conclusion & Next Steps",
])

# 3 ── Problem & Motivation
add_two_col_slide(
    "Problem & Motivation",
    "Advantages",
    [
        "Data efficient — only classifier head is learnable",
        "Fast — single forward pass + seconds-to-minutes training",
        "Interpretable — isolates backbone effect",
    ],
    "Limitations",
    [
        "No task-specific feature adaptation",
        "Resolution sensitivity (pretrained for 224×224)",
        "Linear classifier ceiling",
    ],
    footer="Task: 100-class classification on Mini-ImageNet (50K train / 10K test). "
           "Challenge: small dataset risks overfitting with end-to-end training.",
)

# 4 ── Approach
add_content_slide("Approach: Frozen Backbone + Classical Classifier", [
    "Image → [Pretrained Backbone (frozen)] → Global Avg Pool → Feature Vector → [SVM / LogReg] → Class",
    "",
    "Pipeline:",
    "  1. Load pretrained ImageNet backbone — freeze all weights",
    "  2. Extract feature vectors via forward pass (one-time cost)",
    "  3. Train classical classifier on extracted features",
    "",
    "Classifiers:",
    "  • Linear SVM — maximises geometric margin (hinge loss + L2)",
    "  • Logistic Regression — minimises cross-entropy (calibrated probabilities + L2)",
])

# 5 ── Experimental Design
add_table_slide(
    "Experimental Design — Backbones (3 Families)",
    ["Backbone", "Family", "Params", "Feature Dim", "Key Innovation"],
    [
        ["ConvNeXt-Tiny", "ConvNeXt (2022)", "27.9M", "768", "7×7 DW-Conv, LayerNorm, GELU"],
        ["ResNet-18", "ResNet (2015)", "11.2M", "512", "Basic residual blocks"],
        ["ResNet-34", "ResNet", "21.3M", "512", "Deeper basic blocks"],
        ["ResNet-50", "ResNet", "23.7M", "2048", "Bottleneck blocks"],
        ["EfficientNet-b0", "EfficientNet (2019)", "4.1M", "1280", "MBConv + SE + compound scaling"],
    ],
    intro="32×32: All 5 backbones × {SVM, LogReg}  |  224×224: ConvNeXt-Tiny & ResNet-18 × SVM",
)

# 6 ── Results: Test Accuracy at 32×32
add_table_slide(
    "Results: Test Accuracy at 32×32",
    ["Backbone", "SVM", "LogReg", "Winner"],
    [
        ["ConvNeXt-Tiny", "43.28%", "44.16%", "LogReg (+0.88pp)"],
        ["ResNet-18", "32.52%", "32.46%", "SVM (+0.06pp)"],
        ["ResNet-34", "33.08%", "33.16%", "LogReg (+0.08pp)"],
        ["ResNet-50", "27.34%", "33.22%", "LogReg (+5.88pp)"],
        ["EfficientNet-b0", "23.88%", "24.72%", "LogReg (+0.84pp)"],
    ],
    intro="ConvNeXt-Tiny leads by 10+ pp. LogReg ≥ SVM for all backbones — noisy features favour probabilistic calibration.",
)

# 7 ── Why ConvNeXt-Tiny Dominates
add_two_col_slide(
    "Why ConvNeXt-Tiny Dominates at 32×32",
    "ConvNeXt-Tiny Stem",
    [
        "4×4 stride-4 patchify convolution",
        "32×32 input → 8×8 feature map after stem",
        "7×7 depthwise kernels still have meaningful spatial extent",
        "LayerNorm stabilises activations; GELU preserves gradient flow",
    ],
    "ResNet Stem",
    [
        "7×7 stride-2 conv + 3×3 stride-2 max-pool",
        "32×32 → 8×8 after stem, then quickly 1×1 through stages",
        "Deeper layers receive spatially degenerate feature maps",
        "BatchNorm statistics calibrated for 224×224 distributions",
    ],
    footer="ConvNeXt's patchify stem is resolution-adaptive: it reduces spatial dims in one step "
           "without the cascading downsampling that collapses small inputs in ResNets.",
)

# 8 ── Why LogReg Outperforms SVM
add_table_slide(
    "Why LogReg Outperforms SVM at 32×32",
    ["Property", "SVM (hinge loss)", "LogReg (cross-entropy)"],
    [
        ["Objective", "Maximise geometric margin", "Minimise calibrated log-likelihood"],
        ["Noise handling", "Margin amplifies noisy support vectors", "Probabilistic weighting downweights ambiguous samples"],
        ["High-dim behaviour", "Overfits when signal-to-noise is low", "L2 + cross-entropy regularises more gracefully"],
        ["Solver scaling", "QP: scales poorly with dim (2048-d → 3153s)", "LBFGS: predictable convergence (2048-d → 680s)"],
    ],
    intro="ResNet-50 extreme case: 2048-d noisy features → SVM 27.34% vs LogReg 33.22% (+5.88 pp). "
          "Hinge loss concentrates on noisy support vectors; cross-entropy averages over all samples.",
)

# 9 ── NetScore chart
add_image_slide(
    "Results: NetScore, Accuracy, Inference & Parameters",
    "analysis_netscore_combined.png",
)

# 10 ── NetScore & Efficiency table
add_table_slide(
    "Results: NetScore & Efficiency (32×32)",
    ["Backbone", "Classifier", "Test Acc", "Infer. Time", "Params", "Train Time"],
    [
        ["ConvNeXt-Tiny", "LogReg", "44.16%", "5.92 ms", "27.9M", "278s"],
        ["ConvNeXt-Tiny", "SVM", "43.28%", "5.96 ms", "27.9M", "1502s"],
        ["ResNet-18", "LogReg", "32.46%", "3.09 ms", "11.2M", "347s"],
        ["ResNet-18", "SVM", "32.52%", "2.89 ms", "11.2M", "465s"],
    ],
    intro="NetScore = 20·log₁₀(A² / √T·√P). LogReg trains 1.2–9× faster than SVM. "
          "ResNet-18 has the fastest inference (≈2.9 ms/image).",
)

# 11 ── Efficiency tradeoff chart
add_image_slide(
    "Why Architecture Matters More Than Depth",
    "analysis_efficiency_tradeoff.png",
    caption="ResNets cluster at 32–33% regardless of depth (resolution-limited). "
            "ConvNeXt-Tiny leads by 10+ pp. EfficientNet-b0 collapses.",
)

# 12 ── Why EfficientNet-b0 Fails
add_content_slide("Why EfficientNet-b0 Fails at 32×32", [
    "EfficientNet's compound scaling (Tan & Le, 2019) jointly optimises depth, width, and resolution:",
    "   depth: d = α^φ,   width: w = β^φ,   resolution: r = γ^φ",
    "",
    "Designed for 224×224 (r = 1.0). At 32×32 (r ≈ 0.14):",
    "",
    "1. Squeeze-and-Excitation modules fail — globally average-pool feature maps to compute",
    "   channel attention. When pre-pooling maps are 1×1 or 2×2, the output is random noise.",
    "",
    "2. MBConv depthwise convolutions underperform — 3×3 and 5×5 kernels on 1–2 pixel",
    "   feature maps have no spatial variation to detect.",
    "",
    "3. Depth/width are over-scaled for the resolution — more capacity than the spatial",
    "   information can support, leading to redundant or contradictory features.",
    "",
    "→ Collapsing resolution to 14% while keeping depth/width fixed violates the core design principle.",
])

# 13 ── Resolution Impact chart
add_image_slide(
    "Resolution Impact: 224×224 → 32×32",
    "analysis_accuracy_drop.png",
    caption="Both backbones lose ≈50 pp. ConvNeXt-Tiny's 10.6 pp advantage is preserved across resolutions.",
)

# 14 ── Why ≈50 pp drop
add_content_slide("Why the ≈50 pp Drop Is Nearly Identical", [
    "The uniform ~50 pp drop reveals a shared bottleneck:",
    "",
    "1. Pretrained filters are calibrated for 224×224 — Gabor-like edge detectors in early",
    "   layers expect specific spatial frequencies. At 32×32, the Nyquist frequency is 7×",
    "   lower — fine textures and edges are aliased away before the network sees them.",
    "",
    "2. Receptive field saturation — at 224×224, a ResNet-18 layer-4 neuron covers ~100 px",
    "   (partial image). At 32×32, the same neuron covers the entire image after 2 stages,",
    "   eliminating the hierarchical spatial decomposition that makes CNNs powerful.",
    "",
    "3. The gap is preserved (10.6→10.8 pp) because ConvNeXt's advantage is architectural,",
    "   not resolution-dependent — patchify stem, LayerNorm, and wider kernels extract",
    "   better features at any spatial scale.",
    "",
    "→ ~50 pp is a hard floor for frozen ImageNet backbones on 32×32 — fine-tuning is needed.",
])

# 15 ── t-SNE
add_image_slide(
    "t-SNE Feature Visualisation",
    "downselected_backbones_tsne.png",
    caption="224×224: ConvNeXt-Tiny → tight clusters (93.88%); ResNet-18 → moderate (83.24%). "
            "32×32: ConvNeXt-Tiny retains partial structure; ResNet-18 → near-random.",
)

# 16 ── Generalisation Gap
add_image_slide(
    "Generalisation Gap",
    "analysis_generalisation_gap.png",
    caption="Val − Test gap uniformly small (0.6–3.0 pp). Linear models on frozen features generalise well.",
)

# 17 ── Backbone Downselection
add_two_col_slide(
    "Backbone Downselection for Fine-Tuning",
    "1. ConvNeXt-Tiny (SOTA)",
    [
        "Highest accuracy at all resolutions",
        "93.88% (224×224), 44.16% (32×32)",
        "Modern architecture innovations",
        "Best candidate for further improvement",
    ],
    "2. ResNet-18 (Classical Baseline)",
    [
        "Fastest inference (2.9 ms/image)",
        "Well-established architecture",
        "Measures how much fine-tuning closes the gap",
        "Simple & interpretable",
    ],
    footer="This pairing lets us compare fine-tuning strategies (frozen, partial unfreeze, full, LoRA) "
           "on a modern vs classical architecture.",
)

# 18 ── Downselected comparison chart
add_image_slide(
    "Downselected Backbones — Comparison",
    "downselected_backbones_comparison.png",
)

# 19 ── Key Takeaways
add_content_slide("Key Takeaways", [
    "1. ConvNeXt-Tiny is the best backbone — leads by 10+ pp at any resolution",
    "   thanks to patchify stem and modern CNN design",
    "",
    "2. Resolution is critical — 224→32 causes catastrophic ≈50 pp loss",
    "   for all frozen backbones",
    "",
    "3. LogReg ≥ SVM at 32×32 — noisy features favour probabilistic calibration;",
    "   LogReg also trains 1.2–9× faster",
    "",
    "4. Architecture > Depth — ResNet-18/34/50 cluster within 1 pp;",
    "   design innovations (ConvNeXt) matter far more",
    "",
    "5. EfficientNet-b0 fails at low resolution — compound scaling breaks down",
    "   at 7× below design resolution",
    "",
    "6. Strong generalisation — val-test gap < 3 pp across all configs",
    "   (linear models + frozen features = implicit regularisation)",
])

# 20 ── Next Steps
add_content_slide("Next Steps: Fine-Tuning (Section 3)", [
    "With ConvNeXt-Tiny and ResNet-18 as our downselected backbones, we will explore:",
    "",
    "  • Frozen backbone → linear probe (baseline, already done)",
    "  • Partial unfreezing → unfreeze last N layers",
    "  • Full fine-tuning → all parameters trainable",
    "  • LoRA → low-rank adaptation of backbone weights",
    "",
    "Goal: Determine how much task-specific adaptation can improve over frozen feature",
    "extraction, and whether ConvNeXt-Tiny's architectural advantage persists after fine-tuning.",
])

# 21 ── Thank You
add_title_slide("Thank You", "Questions?  —  CS5242 Project, Traditional ML Approach")

# ── Save ─────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"Exported: {OUT}")
